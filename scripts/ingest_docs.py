#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Ingest PPTX, PDF, HTML, and DOCX files into a JSONL with a unified schema:
{ "doc_id": "...", "source": "...", "where": "...", "type": "pptx|pdf|html|docx", "text": "..." }

Usage example:
  python scripts/ingest_docs.py \
    --pptx_dir data/pptx \
    --pdf_dir  data/pdfs \
    --html_dir data/html \
    --docx_dir data/docx \
    --out_jsonl data/ingested/raw.jsonl
"""
import argparse, os, re, ujson as json
from typing import Iterable

def _norm(s: str) -> str:
    """Collapse whitespace and strip."""
    return re.sub(r'\s+', ' ', (s or '')).strip()

# ---------------- PPTX ----------------
def _extract_pptx(path: str, doc_id: str) -> Iterable[dict]:
    from pptx import Presentation  # python-pptx
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for sh in slide.shapes:
            # text frames (including placeholders)
            if hasattr(sh, "text"):
                t = _norm(getattr(sh, "text", "") or "")
                if t:
                    parts.append(t)
            # tables
            try:
                if getattr(sh, "has_table", False):
                    tbl = sh.table
                    for row in tbl.rows:
                        row_txt = " | ".join(_norm(cell.text) for cell in row.cells)
                        if row_txt:
                            parts.append(row_txt)
            except Exception:
                pass
        text = _norm(" ".join(parts))
        if text:
            yield {
                "doc_id": doc_id,
                "source": os.path.basename(path),
                "where": f"slide {i}",
                "type": "pptx",
                "text": text,
            }

# ---------------- PDF ----------------
def _extract_pdf(path: str, doc_id: str) -> Iterable[dict]:
    from pypdf import PdfReader
    reader = PdfReader(path)
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = _norm(page.extract_text() or "")
        except Exception:
            txt = ""
        if txt:
            yield {
                "doc_id": doc_id,
                "source": os.path.basename(path),
                "where": f"page {i}",
                "type": "pdf",
                "text": txt,
            }

# ---------------- HTML ----------------
def _extract_html(path: str, doc_id: str) -> Iterable[dict]:
    """
    Heuristic HTML → sections:
      - Start a new section at each H1/H2/H3
      - Accumulate subsequent block texts (p, li, td/th, blockquote, pre/code captions)
      - If no headings exist, emit a single 'body' section
    """
    from bs4 import BeautifulSoup  # beautifulsoup4
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "html.parser")

    # Drop non-content tags
    for bad in soup(["script", "style", "noscript", "template"]):
        bad.decompose()

    title = _norm(soup.title.string if soup.title and soup.title.string else os.path.basename(path))
    body = soup.body or soup

    # Collect in reading order
    blocks = body.find_all(
        ["h1","h2","h3","p","li","td","th","blockquote","pre","code","caption","figcaption"],
        recursive=True
    )

    sections = []
    current_heading = None
    current_parts = []

    def flush():
        nonlocal current_parts, current_heading
        text = _norm(" ".join(current_parts))
        if text:
            where = (current_heading or "body")
            sections.append({
                "doc_id": doc_id,
                "source": title,
                "where": where,
                "type": "html",
                "text": text,
            })
        current_parts = []

    for el in blocks:
        name = el.name.lower()
        txt = _norm(el.get_text(separator=" ", strip=True))
        if not txt:
            continue
        if name in ("h1","h2","h3"):
            # New section boundary
            if current_parts:
                flush()
            current_heading = txt
        else:
            current_parts.append(txt)

    # tail
    if current_parts:
        flush()

    # If nothing produced, fall back to full plain text
    if not sections:
        plain = _norm(body.get_text(separator=" ", strip=True))
        if plain:
            sections.append({
                "doc_id": doc_id,
                "source": title,
                "where": "body",
                "type": "html",
                "text": plain,
            })
    return sections

# ---------------- DOCX ----------------
def _extract_docx(path: str, doc_id: str) -> Iterable[dict]:
    """
    DOCX → sections:
      - New section at Heading 1/2/3, Title/Subtitle
      - Accumulate normal paragraphs and table rows under the current section
      - Preserve document order (paragraphs & tables) using oxml walk
    """
    from docx import Document  # python-docx
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    from docx.text.paragraph import Paragraph
    from docx.table import Table

    doc = Document(path)

    def is_heading(para: Paragraph) -> bool:
        try:
            name = (para.style.name or "").lower()
        except Exception:
            name = ""
        return (
            name.startswith("heading")  # Heading 1, Heading 2, ...
            or name in ("title", "subtitle")
        )

    sections = []
    current_heading = None
    current_parts = []

    def flush():
        nonlocal current_parts, current_heading
        text = _norm(" ".join(current_parts))
        if text:
            where = (current_heading or "body")
            sections.append({
                "doc_id": doc_id,
                "source": os.path.basename(path),
                "where": where,
                "type": "docx",
                "text": text,
            })
        current_parts = []

    # Walk in document order (paragraphs + tables)
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            t = _norm(para.text)
            if not t:
                continue
            if is_heading(para):
                if current_parts:
                    flush()
                current_heading = t
            else:
                current_parts.append(t)
        elif isinstance(child, CT_Tbl):
            tbl = Table(child, doc)
            for row in tbl.rows:
                row_txt = " | ".join(_norm(cell.text) for cell in row.cells)
                if row_txt:
                    current_parts.append(row_txt)

    if current_parts:
        flush()

    # If nothing, emit whole doc plain text as single section
    if not sections:
        full = []
        for p in doc.paragraphs:
            tt = _norm(p.text)
            if tt:
                full.append(tt)
        if full:
            sections.append({
                "doc_id": doc_id,
                "source": os.path.basename(path),
                "where": "body",
                "type": "docx",
                "text": _norm(" ".join(full)),
            })
    return sections

# ---------------- CLI ----------------
if __name__ == "__main__":
    ap = argparse.ArgumentParser(description="Ingest PPTX, PDF, HTML, and DOCX into JSONL.")
    ap.add_argument("--pptx_dir", default=None, help="Directory of .pptx files")
    ap.add_argument("--pdf_dir",  default=None, help="Directory of .pdf files")
    ap.add_argument("--html_dir", default=None, help="Directory of .html/.htm files")
    ap.add_argument("--docx_dir", default=None, help="Directory of .docx files")
    ap.add_argument("--out_jsonl", required=True, help="Output JSONL file")
    args = ap.parse_args()

    out_dir = os.path.dirname(args.out_jsonl) or "."
    os.makedirs(out_dir, exist_ok=True)

    n = 0
    with open(args.out_jsonl, "w", encoding="utf-8") as out:
        # PPTX
        if args.pptx_dir and os.path.isdir(args.pptx_dir):
            for fn in sorted(os.listdir(args.pptx_dir)):
                if fn.lower().endswith(".pptx"):
                    full = os.path.join(args.pptx_dir, fn)
                    for rec in _extract_pptx(full, f"pptx::{fn}"):
                        out.write(json.dumps(rec) + "\n"); n += 1
        # PDF
        if args.pdf_dir and os.path.isdir(args.pdf_dir):
            for fn in sorted(os.listdir(args.pdf_dir)):
                if fn.lower().endswith(".pdf"):
                    full = os.path.join(args.pdf_dir, fn)
                    for rec in _extract_pdf(full, f"pdf::{fn}"):
                        out.write(json.dumps(rec) + "\n"); n += 1
        # HTML
        if args.html_dir and os.path.isdir(args.html_dir):
            for fn in sorted(os.listdir(args.html_dir)):
                if fn.lower().endswith((".html", ".htm")):
                    full = os.path.join(args.html_dir, fn)
                    for rec in _extract_html(full, f"html::{fn}"):
                        out.write(json.dumps(rec) + "\n"); n += 1
        # DOCX
        if args.docx_dir and os.path.isdir(args.docx_dir):
            for fn in sorted(os.listdir(args.docx_dir)):
                if fn.lower().endswith(".docx"):
                    full = os.path.join(args.docx_dir, fn)
                    for rec in _extract_docx(full, f"docx::{fn}"):
                        out.write(json.dumps(rec) + "\n"); n += 1

    print(f"Wrote {n} records -> {args.out_jsonl}")