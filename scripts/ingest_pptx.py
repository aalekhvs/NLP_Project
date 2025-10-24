import argparse, json, os
from pathlib import Path
from pptx import Presentation

def extract_pptx(pptx_path):
    pres = Presentation(pptx_path)
    records = []
    for i, slide in enumerate(pres.slides, 1):
        title = ""
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if "title" in shape.name.lower():
                    title = (shape.text or "").strip()
                    break
        body_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                t = (shape.text or "").strip()
                if t: body_texts.append(t)
            if hasattr(shape, "has_table") and shape.has_table:
                rows = []
                for r in shape.table.rows:
                    rows.append(" | ".join([c.text.strip() for c in r.cells]))
                if rows: body_texts.append("\n".join(rows))
        notes = ""
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        full_text = "\n".join([t for t in [title] + body_texts + [notes] if t])
        if full_text.strip():
            records.append({
                "doc_id": Path(pptx_path).stem,
                "slide_number": i,
                "slide_title": title,
                "text": full_text,
                "version": "v1"
            })
    return records

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx_dir", required=True)
    ap.add_argument("--out", required=True)
    args = ap.parse_args()

    out = Path(args.out); out.parent.mkdir(parents=True, exist_ok=True)
    count = 0
    with out.open("w", encoding="utf-8") as f:
        for fname in os.listdir(args.pptx_dir):
            if fname.lower().endswith(".pptx"):
                recs = extract_pptx(os.path.join(args.pptx_dir, fname))
                for r in recs:
                    f.write(json.dumps(r, ensure_ascii=False) + "\n")
                    count += 1
    print(f"Wrote {count} slide records → {out}")

if __name__ == "__main__":
    main()