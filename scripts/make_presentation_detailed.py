#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate a dense, research-style PPT (~36 slides) for the Course FAQ RAG project.
- Lots of text (multi-level bullets) + speaker notes so prof has fewer follow-ups.
- Charts via matplotlib only (one chart per figure, default colors).
- Optional case studies: reads top hits from data/eval/retrievals_reranked.jsonl (else retrievals.jsonl).
"""

import os, json, itertools
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# ------------------ METRICS FROM YOUR RUNS ------------------
baseline = {
    "Queries": 215,
    "MRR": 0.3022,
    "Hit@1": 0.2140,
    "Hit@3": 0.3488,
    "Hit@5": 0.4279,
    "Hit@10": 0.5628,
    "nDCG@1": 0.2140,
    "nDCG@3": 0.2736,
    "nDCG@5": 0.3248,
    "nDCG@10": 0.3053,
    "EM": 0.0884,
    "F1": 0.2532,
}
fusion = {
    "Queries": 215,
    "MRR": 0.3252,
    "Hit@1": 0.2372,
    "Hit@3": 0.3395,
    "Hit@5": 0.4558,
    "Hit@10": 0.5628,
    "nDCG@1": 0.2372,
    "nDCG@3": 0.2981,
    "nDCG@5": 0.3449,
    "nDCG@10": 0.3814,
    "EM": 0.2000,   # best normalized EM you reported for fusion-only (top-1)
    "F1": 0.3404,
}
final_res = {
    "Queries": 215,
    "MRR": 0.4319,
    "Hit@1": 0.3349,
    "Hit@3": 0.4837,
    "Hit@5": 0.5721,
    "Hit@10": 0.6605,
    "nDCG@1": 0.3349,
    "nDCG@3": 0.4215,
    "nDCG@5": 0.4577,
    "nDCG@10": 0.4864,
    "EM": 0.3256,
    "F1": 0.4451,
}

FIG_DIR = "figs"
os.makedirs(FIG_DIR, exist_ok=True)

# ------------------ FIGURE HELPERS ------------------
def bar_fig(labels, vals, title, out_png):
    plt.figure()
    plt.bar(labels, vals)  # default colors, one chart only
    plt.ylim(0, 1.0)
    plt.ylabel("Score")
    plt.title(title)
    plt.tight_layout()
    plt.savefig(out_png, dpi=200)
    plt.close()

def line_fig(labels, vals, title, out_png):
    plt.figure()
    plt.plot(labels, vals, marker='o')
    plt.ylim(0, 1.0)
    plt.ylabel("Score")
    plt.title(title)
    plt.tight_layout()
    plt.savefig(out_png, dpi=200)
    plt.close()

def pipeline_fig(out_png):
    plt.figure(figsize=(11, 3))
    ax = plt.gca()
    ax.axis('off')
    boxes = [
        ("Ingest\nPDF/PPTX/HTML/DOCX", 0.06, 0.5),
        ("Chunk\n(180/30)", 0.24, 0.5),
        ("Index\nTF–IDF", 0.38, 0.5),
        ("Retrieve\n+ BM25 + RRF\n+ Category + Heading", 0.56, 0.5),
        ("Rerank\nCross-Encoder", 0.74, 0.5),
        ("Extract + Normalize\n(EM/F1 eval)", 0.90, 0.5),
    ]
    for text, x, y in boxes:
        ax.add_patch(plt.Rectangle((x-0.07, y-0.12), 0.14, 0.24, fill=False))
        ax.text(x, y, text, ha='center', va='center')
    for i in range(len(boxes)-1):
        x0 = boxes[i][1] + 0.07
        x1 = boxes[i+1][1] - 0.07
        y = boxes[i][2]
        ax.annotate("", xy=(x1, y), xytext=(x0, y), arrowprops=dict(arrowstyle="->"))
    plt.tight_layout()
    plt.savefig(out_png, dpi=220, bbox_inches='tight')
    plt.close()

# Generate figures
bar_fig(["Hit@1","Hit@3","Hit@5","Hit@10"], [baseline[k] for k in ["Hit@1","Hit@3","Hit@5","Hit@10"]],
        "Baseline (TF–IDF) — Hit Rates", f"{FIG_DIR}/baseline_hit.png")
bar_fig(["nDCG@1","nDCG@3","nDCG@5","nDCG@10"], [baseline[k] for k in ["nDCG@1","nDCG@3","nDCG@5","nDCG@10"]],
        "Baseline (TF–IDF) — nDCG", f"{FIG_DIR}/baseline_ndcg.png")

bar_fig(["Hit@1","Hit@3","Hit@5","Hit@10"], [fusion[k] for k in ["Hit@1","Hit@3","Hit@5","Hit@10"]],
        "Fusion (BM25+RRF+filters) — Hit Rates", f"{FIG_DIR}/fusion_hit.png")
bar_fig(["nDCG@1","nDCG@3","nDCG@5","nDCG@10"], [fusion[k] for k in ["nDCG@1","nDCG@3","nDCG@5","nDCG@10"]],
        "Fusion (BM25+RRF+filters) — nDCG", f"{FIG_DIR}/fusion_ndcg.png")

bar_fig(["Hit@1","Hit@3","Hit@5","Hit@10"], [final_res[k] for k in ["Hit@1","Hit@3","Hit@5","Hit@10"]],
        "Final (Fusion + CE) — Hit Rates", f"{FIG_DIR}/final_hit.png")
bar_fig(["nDCG@1","nDCG@3","nDCG@5","nDCG@10"], [final_res[k] for k in ["nDCG@1","nDCG@3","nDCG@5","nDCG@10"]],
        "Final (Fusion + CE) — nDCG", f"{FIG_DIR}/final_ndcg.png")

line_fig(["Baseline","Fusion","Final"], [baseline["Hit@1"], fusion["Hit@1"], final_res["Hit@1"]],
         "Progress: Hit@1 (Baseline → Fusion → Final)", f"{FIG_DIR}/progress_hit1.png")
line_fig(["Baseline","Fusion","Final"], [baseline["MRR"], fusion["MRR"], final_res["MRR"]],
         "Progress: MRR (Baseline → Fusion → Final)", f"{FIG_DIR}/progress_mrr.png")

bar_fig(["Baseline","Fusion","Final"], [baseline["EM"], fusion["EM"], final_res["EM"]], "QA EM (normalized)", f"{FIG_DIR}/qa_em.png")
bar_fig(["Baseline","Fusion","Final"], [baseline["F1"], fusion["F1"], final_res["F1"]], "QA F1 (normalized)", f"{FIG_DIR}/qa_f1.png")

pipeline_fig(f"{FIG_DIR}/pipeline.png")

# ------------------ PPT HELPERS ------------------
def add_title(prs, title, subtitle=None):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    if subtitle:
        s.placeholders[1].text = subtitle
    return s

def add_bullets(prs, title, bullets, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = title
    tf = s.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(b, (list, tuple)):
            # first item is a parent bullet; rest are level-1 children
            parent = b[0]
            p.text = parent
            p.level = 0
            for sub in b[1:]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.level = 1
        else:
            p.text = b
            p.level = 0
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s

def add_picture(prs, title, img_path, notes=None, height_in=4.2):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    s.shapes.title.text = title
    left, top = Inches(0.75), Inches(1.5)
    s.shapes.add_picture(img_path, left, top, height=Inches(height_in))
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s

def add_two_pictures(prs, title, img1, img2, notes=None, height_in=3.8):
    s = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    s.shapes.title.text = title
    left1, left2, top = Inches(0.5), Inches(5.1), Inches(1.6)
    s.shapes.add_picture(img1, left1, top, height=Inches(height_in))
    s.shapes.add_picture(img2, left2, top, height=Inches(height_in))
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s

def add_table(prs, title, rows, notes=None):
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = title
    cols = len(rows[0])
    table = s.shapes.add_table(rows=len(rows), cols=cols,
                               left=Inches(0.5), top=Inches(1.6),
                               width=Inches(9.3), height=Inches(0.9*len(rows))).table
    for r in range(len(rows)):
        for c in range(cols):
            table.cell(r, c).text = str(rows[r][c])
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s

# ------------------ OPTIONAL CASE STUDIES ------------------
def maybe_add_case_studies(prs, chunks_path="data/chunks/chunks.jsonl"):
    # pick retrievals file if present
    retr_path = None
    for p in ["data/eval/retrievals_reranked.jsonl", "data/eval/retrievals.jsonl"]:
        if os.path.exists(p):
            retr_path = p
            break
    if not retr_path or not os.path.exists(chunks_path):
        return

    # load first 3 queries with hits
    with open(retr_path, "r", encoding="utf-8") as f:
        entries = []
        for ln in itertools.islice(f, 6):  # read a few; we’ll filter
            if ln.strip():
                entries.append(json.loads(ln))
    entries = [e for e in entries if e.get("hits")]
    entries = entries[:3] if len(entries) >= 3 else entries

    # load chunk texts
    chunk_text = {}
    with open(chunks_path, "r", encoding="utf-8") as f:
        for ln in f:
            if ln.strip():
                rec = json.loads(ln)
                chunk_text[rec["chunk_id"]] = rec.get("text", "")

    for i, e in enumerate(entries, 1):
        q = e["question"]
        hits = e["hits"][:3]
        bullets = [
            ["Query", q],
            "Top hits (IDs, why-keys, first ~140 chars):"
        ]
        for h in hits:
            cid = h["chunk_id"]
            why = ", ".join([f"{k}={v:.3f}" for k,v in h.get("why",{}).items() if isinstance(v,(int,float))])
            txt = (chunk_text.get(cid,"")[:140] + " …") if chunk_text.get(cid) else "(text not found)"
            bullets.append([f"{cid}", f"Scores: {why}", f"Snippet: {txt}"])
        add_bullets(prs, f"Case Study {i}: Retrieval Inspection", bullets,
                    notes="Shows real query and why top hits ranked high; use to preempt probing about ranking rationale.")

# ------------------ BUILD THE DETAILED DECK ------------------
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)

# 1 Title
add_title(prs, "Course FAQ RAG: From Sparse Baseline to Fusion + CE Rerank",
          "Aalekh Srivastava • CS593 NLP • Final Presentation")

# 2 Abstract (dense)
add_bullets(prs, "Abstract", [
    ["Objective", "Build a reliable, updatable RAG assistant to answer course FAQs grounded in course materials with citations."],
    ["Approach Summary",
     "Sparse baseline (TF–IDF) for robustness on small, domain-specific corpora.",
     "Heterogeneous retrieval via BM25 + Reciprocal Rank Fusion (RRF) to improve recall under phrasing variation.",
     "Category-aware and heading-aware boosts to steer admin/policy queries to appropriate documents.",
     "Cross-encoder reranking for first-hit precision, followed by extractive answer with normalization."],
    ["Outcomes",
     f"Retrieval improved from MRR {baseline['MRR']:.3f} → {final_res['MRR']:.3f}, Hit@1 {baseline['Hit@1']:.3f} → {final_res['Hit@1']:.3f}.",
     f"QA (normalized): EM {baseline['EM']:.3f} → {final_res['EM']:.3f}, F1 {baseline['F1']:.3f} → {final_res['F1']:.3f}."],
], notes="Keep this slide tight but comprehensive; answers the 'what/why/how/outcome' in one place.")

# 3 Problem Statement & Research Questions
add_bullets(prs, "Problem Statement & Research Questions", [
    ["Problem", "Students repeatedly ask admin/policy/HW logistics that exist in course materials; answers must be exact and citable."],
    ["RQs",
     "RQ1: On a small, curated course corpus, how far can sparse retrieval go before needing heavy neural indexing?",
     "RQ2: Does BM25+RRF fusion and lightweight metadata boosts raise first-hit accuracy meaningfully?",
     "RQ3: What is the marginal benefit of a compact cross-encoder reranker on top-N candidates?"],
    ["Constraints", "Low ops overhead, reproducibility, fast updates when materials change."],
])

# 4 Motivation & Impact (dense)
add_bullets(prs, "Motivation & Impact", [
    ["For instructors", "Lower repeated Q&A load; consistent policy answers; easy to update by re-running pipeline."],
    ["For students", "Fast, precise responses; citations for trust and self-serve verification."],
    ["For research", "Evaluates classical IR + light ML stacking for a realistic, small-corpus RAG setting."],
])

# 5 Related Work: Retrieval & Fusion
add_bullets(prs, "Related Work: Retrieval & Fusion", [
    ["Sparse Retrieval", "TF–IDF/BM25 strong on exact/near-lexical matches; robust with small corpora."],
    ["Fusion", "Reciprocal Rank Fusion combines diverse signals without calibration; widely used in IR competitions."],
    ["Applicability", "Our corpus is small and structured (syllabus/HW/lectures); sparse + fusion is a strong first choice."],
])

# 6 Related Work: Reranking & QA
add_bullets(prs, "Related Work: Reranking & Extractive QA", [
    ["Reranking", "Cross-encoders (e.g., MS MARCO MiniLM-L-6-v2) give strong pairwise scoring at moderate cost."],
    ["Extractive QA", "Stable for policy-style prose; avoids hallucination; normalization improves EM comparability."],
])

# 7 Data Sources & Ingestion
add_bullets(prs, "Data Sources & Ingestion", [
    ["Sources", "PDF, PPTX, HTML, DOCX"],
    ["Parsing", "Preserve doc_id, type, and 'where' (e.g., heading vs body) to enable metadata-aware boosts."],
    ["Schema", "Unified JSONL: {doc_id, type, where, text}"],
], notes="Emphasize reproducibility and metadata preservation: it enables heading/category strategies later.")

# 8 Chunking Strategy
add_bullets(prs, "Chunking Strategy", [
    ["Default", "180 words with 30-word overlap"],
    ["Why", "Short enough to isolate policy sentences; overlap reduces boundary misses."],
    ["Ablation", "240/60 diluted key policy cues; reduced Hit@1 and MRR."],
])

# 9 Gold Construction
add_bullets(prs, "Gold Construction", [
    ["Size & Coverage", "215 curated Q–A with gold chunk IDs; admin/policy, HW logistics, topic concepts."],
    ["Method", "Seeded by course artifacts; auto-gen drafts + manual curation for correctness and coverage."],
    ["Plan", "Scale to 400–600 for tighter confidence intervals and slice-level ablations."],
])

# 10 Pipeline Diagram
add_picture(prs, "System Overview", f"{FIG_DIR}/pipeline.png",
            notes="Ingest → Chunk → Index → Retrieve (TF–IDF + BM25 + RRF + boosts) → Rerank (CE) → Extract + Normalize → Evaluate.")

# 11 Methodology: Index & Baseline
add_bullets(prs, "Methodology: Index & TF–IDF Baseline", [
    ["Index", "Sklearn TF–IDF; cosine similarity over chunk vectors."],
    ["Baseline Use", "Establishes robust lexical baseline; fast to rebuild when materials change."],
])

# 12 Methodology: BM25 & RRF
add_bullets(prs, "Methodology: BM25 & RRF", [
    ["BM25", "Compliments TF–IDF by emphasizing term frequency and doc length normalization."],
    ["RRF", "Fuse TF–IDF and BM25 ranks; score = Σ 1/(K + rank); K=60."],
    ["Effect", "Improves recall under phrasing/synonym changes without extra training."],
])

# 13 Methodology: Category-Aware Boosts
add_bullets(prs, "Methodology: Category-Aware Boosts", [
    ["Heuristic", "Classify query as policy / HW / topics via keywords."],
    ["Soft Filter", "Multiply scores for allowed docs (e.g., syllabus/policies for policy queries)."],
    ["Why", "Removes distractors early; aligns intent with likely sources."],
])

# 14 Methodology: Heading-Aware Boosts
add_bullets(prs, "Methodology: Heading-Aware Boosts", [
    ["Signal", "If chunk 'where' looks like a section heading (e.g., 'Course Policies'), apply multiplicative boost."],
    ["Why", "Policies and rubrics are often summarized in headings; promotes precision for rank-1."],
])

# 15 Methodology: CE Reranker
add_bullets(prs, "Methodology: Cross-Encoder Reranker", [
    ["Model", "cross-encoder/ms-marco-MiniLM-L-6-v2 (compact, proven on MS MARCO)."],
    ["Usage", "Rerank top-N candidates; increases first-hit precision significantly."],
    ["Fallback", "If CE unavailable offline, fall back to token-overlap rerank to retain pipeline operability."],
])

# 16 Methodology: Answer Extraction & Normalization
add_bullets(prs, "Methodology: Answer Extraction & Normalization", [
    ["Extraction", "Select best-matching sentence from top-k concatenated context by token-overlap scoring."],
    ["Normalization", "Canonicalize dates/percents/symbols to improve EM fairness and comparability."],
])

# 17 Experimental Setup (dense)
add_bullets(prs, "Experimental Setup", [
    ["Data", f"215 queries; chunking 180/30; multi-format corpus (PDF/PPTX/HTML/DOCX)."],
    ["Params", "k ∈ {15,40}; topN=150; RRF K=60; category soft boost=0.4; heading weight=0.75."],
    ["Metrics", "Retrieval: MRR, Hit@k, nDCG@k. QA: EM, F1 (normalized)."],
    ["Repro", "run_pipeline.py builds indices, runs retrieval/rerank, evaluates, and writes tables."],
])

# 18–21 Baseline/Fusion/Final charts
add_two_pictures(prs, "Results: Baseline (TF–IDF)", f"{FIG_DIR}/baseline_hit.png", f"{FIG_DIR}/baseline_ndcg.png")
add_two_pictures(prs, "Results: Fusion (BM25+RRF+filters)", f"{FIG_DIR}/fusion_hit.png", f"{FIG_DIR}/fusion_ndcg.png")
add_two_pictures(prs, "Results: Final (Fusion + CE)", f"{FIG_DIR}/final_hit.png", f"{FIG_DIR}/final_ndcg.png")
add_two_pictures(prs, "Retrieval Progress", f"{FIG_DIR}/progress_hit1.png", f"{FIG_DIR}/progress_mrr.png")

# 22–23 QA bars
add_two_pictures(prs, "QA (normalized): EM and F1", f"{FIG_DIR}/qa_em.png", f"{FIG_DIR}/qa_f1.png",
                 notes="Normalization mattered for EM; F1 also rose with CE rerank thanks to better first-hit context.")

# 24 Ablation Table
rows = [
    ["System", "MRR", "Hit@1", "Hit@5", "EM", "F1"],
    ["TF–IDF", f"{baseline['MRR']:.3f}", f"{baseline['Hit@1']:.3f}", f"{baseline['Hit@5']:.3f}", f"{baseline['EM']:.3f}", f"{baseline['F1']:.3f}"],
    ["+ BM25 + RRF + filters", f"{fusion['MRR']:.3f}", f"{fusion['Hit@1']:.3f}", f"{fusion['Hit@5']:.3f}", f"{fusion['EM']:.3f}", f"{fusion['F1']:.3f}"],
    ["+ CE rerank + normalize", f"{final_res['MRR']:.3f}", f"{final_res['Hit@1']:.3f}", f"{final_res['Hit@5']:.3f}", f"{final_res['EM']:.3f}", f"{final_res['F1']:.3f}"],
]
add_table(prs, "Ablation Summary (Key Metrics)", rows,
          notes="Call out the *largest marginal gain* from CE on Hit@1/MRR; and normalization elevating EM.")

# 25 Error Analysis (file/type slices – narrative)
add_bullets(prs, "Error Analysis (Slices)", [
    ["Policy vs Topics", "Policy queries occasionally tied between syllabus and lecture mention; heading boost helped but ties remain."],
    ["PPTX vs PDF", "Topic definitions in slides may be brief; extractive single-sentence can miss adjoining clarifications."],
    ["DOCX (HW)", "HW logistics spread across bullet lists; chunking preserved most, but sometimes boundary misses occur."],
])

# 26 What Helped Most (dense narrative)
add_bullets(prs, "What Helped Most", [
    ["Fusion (BM25+RRF)", "Recovered lexical variations; lifted tail queries that TF–IDF alone under-ranked."],
    ["Category + Heading", "Steered admin/policy questions directly to relevant sections; noticeable Hit@1 gains."],
    ["CE Rerank", "Largest single jump in first-hit precision; especially for phrasing-mismatch queries."],
])

# 27 Ideas Tried That Didn’t Work (yet)
add_bullets(prs, "Ideas Tried That Didn’t Work (Yet)", [
    ["Larger chunks (240/60)", "Diluted policy sentences; lower Hit@1 and MRR."],
    ["Naive synonym expansion", "Introduced noise; BM25 already covers many variants; requires calibrated query rewriting/fusion weights."],
])

# 28 Compute/Latency & Ops
add_bullets(prs, "Compute, Latency & Ops", [
    ["Indexing", "Sparse vectorizer rebuild is fast; cheap to refresh when materials change."],
    ["Rerank Cost", "CE applied only to top-N; tradeoff tunable (we used MiniLM-L-6-v2)."],
    ["Ops", "Single script (`run_pipeline.py`) ends with CSV tables and report-ready MD."],
])

# 29 Threats to Validity
add_bullets(prs, "Threats to Validity", [
    ["Gold Size", "215 Q–A is modest; confidence intervals wide for some slices."],
    ["Drift", "Policies change over semesters; requires rebuild and partial re-curation."],
    ["Bias", "Heuristic category detection may misclassify edge queries."],
])

# 30 Research Framing: Why This Design
add_bullets(prs, "Research Framing: Why This Design", [
    ["Small-Corpus Reality", "Sparse + light ML stacking is a pragmatic sweet spot vs heavy dense retrieval."],
    ["Explainability", "Sparse + CE rerank with extractive answers is inspectable and auditable by instructors."],
    ["Extensibility", "Each stage (fusion, boosts, rerank, normalization) can be ablated and improved independently."],
])

# 31 Case Studies (optional, auto from your files if present)
maybe_add_case_studies(prs)

# 32 Progress Since Update 1 (dense)
add_bullets(prs, "Progress Since Update 1", [
    ["Engineering", "DOCX/HTML ingestion added; preserved headings/metadata; scripted E2E pipeline with reproducible outputs."],
    ["Evaluation", "Expanded gold to 215; added normalized QA evaluation; charts + ablation tables auto-generated."],
    ["Results", f"MRR {baseline['MRR']:.3f}→{final_res['MRR']:.3f}; Hit@1 {baseline['Hit@1']:.3f}→{final_res['Hit@1']:.3f}; EM {baseline['EM']:.3f}→{final_res['EM']:.3f}."],
])

# 33 Limitations
add_bullets(prs, "Limitations", [
    ["Extractive scope", "Single-sentence extraction may underserve multi-sentence definitions/rubrics."],
    ["Heuristics", "Category/heading boosts are simple; could be learned or replaced by structured indexing."],
    ["Scale", "Gold set needs to grow (400–600) for robust slice-wise ablations."],
])

# 34 Next Steps (planned work)
add_bullets(prs, "Next Steps (Planned)", [
    ["Gold Expansion", "400–600 Q–A; stratified by category/filetype."],
    ["Learned Signals", "Lightweight learned rerank/fusion, calibrated cutoffs."],
    ["Normalization+", "Richer canonicalization (dates, ranges, penalties) to raise EM further."],
])

# 35 Conclusion (dense)
add_bullets(prs, "Conclusion", [
    ["Deliverable", "A reproducible, updatable course-FAQ RAG with strong retrieval gains and clean ablation story."],
    ["Takeaway", "Classical IR + light ML stacking is highly competitive for small curated corpora, with clear knobs for precision."],
    ["Readiness", "Pipeline + charts + tables are presentation- and paper-ready; extensible for future semesters."],
])

# 36 Backup: Metrics Definitions (text)
add_bullets(prs, "Appendix: Metrics Definitions (Plain Text)", [
    ["MRR", "Mean Reciprocal Rank; average of 1/rank of the first relevant result across queries."],
    ["Hit@k", "Fraction of queries with at least one relevant result in top-k."],
    ["nDCG@k", "Normalized Discounted Cumulative Gain at k; higher weight for early relevant ranks."],
    ["EM/F1", "Exact string match after normalization; token-level F1 between prediction and gold."],
])

OUT = "NLP_Project_Presentation_Aalekh_Final_DETAILED.pptx"
prs.save(OUT)
print(f"Saved -> {OUT}")